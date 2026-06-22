# app/documents/parsers/pptx_parser.py
#
# PPTX extraction via python-pptx. One Document per slide, including
# speaker notes. Tables converted to markdown. Content-bearing images
# (charts, diagrams, screenshots) are described via vision LLM and
# inserted into the slide's text; decorative/repeated images (logos,
# icons) are filtered out before any LLM call.
#
# Verified API used here (current python-pptx 1.0.0 docs):
#   has_text_frame / text_frame.text
#   has_table / table
#   has_notes_slide / notes_slide.notes_text_frame.text
#   shape.shape_type == MSO_SHAPE_TYPE.PICTURE
#   shape.image.blob — documented, public attribute
#
# IMPORTANT CORRECTION: shape.image.size does NOT exist on python-pptx's
# Image object. The only documented Image properties are blob,
# content_type, dpi, ext, filename, sha1 — no size/dimensions property.
# Pixel dimensions are obtained here via Pillow (PIL), which is already
# an installed transitive dependency of python-pptx itself (confirmed:
# python-pptx's own internal ImagePart class uses PIL internally to
# read image dimensions) — so this adds no new dependency, it just uses
# the public, documented blob attribute plus Pillow directly, rather
# than relying on python-pptx's private/undocumented internals.

import logging
from io import BytesIO

from langchain_core.documents import Document
from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.documents.image_description import describe_image
from app.documents.image_filter import reset_seen_images, should_describe_image

logger = logging.getLogger("app.documents.parsers.pptx")


async def parse_pptx(file_bytes: bytes, filename: str) -> list[Document]:
    """One Document per slide."""
    prs = Presentation(BytesIO(file_bytes))
    reset_seen_images()
    documents = []

    for i, slide in enumerate(prs.slides):
        text_parts: list[str] = []
        described_count = 0
        skipped_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_parts.append(shape.text_frame.text)
            elif shape.has_table:
                text_parts.append(_table_to_markdown(shape.table))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                described, skipped = await _process_picture_shape(
                    shape, filename, i + 1, text_parts
                )
                described_count += described
                skipped_count += skipped

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                text_parts.append(f"[Speaker notes: {notes_text}]")

        if described_count or skipped_count:
            logger.info(
                "%s slide %d: %d image(s) described, %d skipped "
                "(decorative/repeated/failed)",
                filename, i + 1, described_count, skipped_count,
            )

        documents.append(
            Document(
                page_content="\n\n".join(text_parts),
                metadata={
                    "source": filename,
                    "file_type": "pptx",
                    "slide": i + 1,
                    "images_described": described_count,
                    "images_skipped": skipped_count,
                },
            )
        )
    return documents


async def _process_picture_shape(
    shape, filename: str, slide_num: int, text_parts: list[str]
) -> tuple[int, int]:
    """
    Extracts a picture shape's image bytes, determines real pixel
    dimensions via Pillow, decides whether it's content-bearing,
    and if so describes it via vision LLM and appends to text_parts.

    Returns (described_count, skipped_count) — each 0 or 1, since
    this handles exactly one shape.
    """
    try:
        image_bytes = shape.image.blob

        with PILImage.open(BytesIO(image_bytes)) as pil_img:
            width, height = pil_img.size

        if should_describe_image(image_bytes, width, height):
            description = await describe_image(image_bytes)
            text_parts.append(f"[Image description: {description}]")
            return 1, 0
        else:
            return 0, 1

    except Exception as e:
        logger.warning(
            "%s slide %d: image processing failed: %s",
            filename, slide_num, e,
        )
        return 0, 1


def _table_to_markdown(table) -> str:
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body = "\n".join("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n".join([header, separator, body])